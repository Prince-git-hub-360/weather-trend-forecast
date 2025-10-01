# ============================================================
# Weather Trend Forecasting: Generate Presentation
# ============================================================

# Step 0: Imports & Setup
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ---------------------------
# Step 1: Project Paths
# ---------------------------
project_root = os.path.abspath(os.path.join(os.getcwd(), ".."))

# Use existing 'reports' folder
reports_dir = os.path.join(project_root, "reports")
# Make sure the folder exists
os.makedirs(reports_dir, exist_ok=True)

# Output PPTX file
output_pptx = os.path.join(reports_dir, "Weather_Trend_Forecasting.pptx")

# ---------------------------
# Step 2: Load cleaned data
# ---------------------------
cleaned_csv_path = os.path.join(project_root, "weather-trend-forecast","data", "processed", "weather_cleaned.csv")
df = pd.read_csv(cleaned_csv_path)
df["last_updated"] = pd.to_datetime(df["last_updated"])

# ---------------------------
# Step 3: Create Presentation
# ---------------------------
prs = Presentation()

# ---------------------------
# Slide 1: Title Slide
# ---------------------------
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Weather Trend Forecasting"
slide.placeholders[1].text = "PM Accelerator | AI Project | Global Weather Repository"

# ---------------------------
# Slide 2: PM Accelerator Mission
# ---------------------------
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "PM Accelerator Mission"
mission_text = (
    "I’m on a mission to help launch 1,000+ AI products and empower professionals like you "
    "to become the next generation of AI product leaders — impacting millions of lives "
    "through real-world innovation."
)
slide.placeholders[1].text = mission_text

# ---------------------------
# Slide 3: Data Overview
# ---------------------------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Data Overview"
overview_text = (
    f"Total rows: {len(df)}\n"
    f"Total locations: {df['location_name'].nunique()}\n"
    f"Total countries: {df['country'].nunique()}\n"
    f"Date range: {df['last_updated'].min().date()} to {df['last_updated'].max().date()}"
)
slide.placeholders[1].text = overview_text

# ---------------------------
# Slide 4: EDA Insights
# ---------------------------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "EDA Insights"
eda_summary = df[['temperature_celsius','precip_mm','humidity','wind_kph']].describe().round(2).to_string()
slide.placeholders[1].text = f"Descriptive statistics:\n{eda_summary}"

# ---------------------------
# Slide 5: Forecasting Summary
# ---------------------------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Forecasting Summary"
forecast_text = (
    "RandomForest model was trained to predict temperature based on humidity, precipitation, "
    "wind, and cloud coverage.\n"
    "MAE & RMSE metrics computed to evaluate model performance.\n"
    "Feature importance and actual vs predicted plots generated."
)
slide.placeholders[1].text = forecast_text

# ---------------------------
# Slide 6: Advanced Analysis
# ---------------------------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Advanced Analysis"
adv_text = (
    "Anomaly detection using IsolationForest identified extreme temperature events.\n"
    "Seasonal trends by month were analyzed.\n"
    "Top countries with highest average temperatures were visualized.\n"
    "Spatial patterns and climate trends explored across locations."
)
slide.placeholders[1].text = adv_text

# ---------------------------
# Slide 7: Next Steps
# ---------------------------
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Next Steps / Recommendations"
next_steps = (
    "- Extend forecasting to all major locations\n"
    "- Build ensemble models to improve forecast accuracy\n"
    "- Explore climate trends across countries and continents\n"
    "- Visualize spatial patterns using GIS maps"
)
slide.placeholders[1].text = next_steps

# ---------------------------
# Step 4: Save Presentation
# ---------------------------
prs.save(output_pptx)
print(f"✅ Presentation saved to: {output_pptx}")
